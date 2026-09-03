#!/usr/bin/env -S uv run --script
# /// script
# requires-python = ">=3.11"
# dependencies = [
#   "python-pptx>=1.0.2",
# ]
# ///
"""Regenerate the fixture decks under ``fixtures/``.

Two pairs of one-slide decks.

``ref.pptx`` / ``gen.pptx`` are identical apart from three deliberately planted
discrepancies, so ``test_compare_decks.py`` can assert the parse finds exactly
those and nothing else:

1. the title is 28pt in ref and 24pt in gen — a formatting delta;
2. the rectangle sits 0.10in further right in gen — a coordinate delta, chosen
   to fall outside the 0.04in default tolerance but inside a 0.25in one, so the
   same pair of fixtures exercises the ``--pos-tol`` flag both ways;
3. gen carries an extra textbox that ref does not have — an unmatched shape.

``coalesced-ref.pptx`` / ``coalesced-gen.pptx`` are the fold-flag pair: the same
paragraph written as one run on the ref side and as three identically formatted
runs on the gen side, which is what PowerPoint's run coalescing does to a deck
pptxgenjs emitted. Everything else about the two is equal, so the comparison has
exactly one thing to say about them and it is an engine artefact — the pair a
``--fold-engine-artefacts`` run has to reduce to zero.

All four decks are built with python-pptx rather than the JS generator the sibling
``pptx-arch-style`` skill uses for its fixtures: python-pptx is already a dev
dependency of this repository, so the tests need neither Node nor a vendored
``node_modules``.

Usage:
    uv run gen_fixtures.py
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches, Pt

HERE = Path(__file__).resolve().parent
FIXTURES = HERE / "fixtures"

TITLE_TEXT = "Structural comparison"
RECT_X_IN = 0.60
GEN_RECT_X_IN = 0.70  # planted delta 2: 0.10in right of ref

# The one paragraph of the coalesced pair, split at the spaces on the gen side.
# The fragments concatenate to the ref text exactly, so the two paragraphs are
# equal as text and differ only in how many <a:r> elements carry it.
COALESCED_TEXT = "Alpha beta gamma"
SPLIT_TEXT = ("Alpha ", "beta ", "gamma")


def build(path: Path, *, title_size_pt: int, rect_x_in: float, extra: bool) -> None:
    """Write a one-slide deck with the given planted variations."""
    presentation = Presentation()
    presentation.slide_width = Inches(10)
    presentation.slide_height = Inches(5.625)
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])

    title = slide.shapes.add_textbox(
        Inches(0.75), Inches(0.20), Inches(8.5), Inches(0.6)
    )
    title.name = "title"
    run = title.text_frame.paragraphs[0].add_run()
    run.text = TITLE_TEXT
    run.font.name = "Arial"
    run.font.size = Pt(title_size_pt)
    run.font.bold = True

    rect = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(rect_x_in), Inches(1.20), Inches(3.0), Inches(1.0)
    )
    rect.name = "body-box"

    if extra:
        box = slide.shapes.add_textbox(
            Inches(6.00), Inches(1.20), Inches(3.0), Inches(0.5)
        )
        box.name = "extra-note"
        box.text_frame.paragraphs[0].add_run().text = "only in gen"

    presentation.save(str(path))


def build_run_split(path: Path, fragments: tuple[str, ...]) -> None:
    """Write a one-slide deck holding ``fragments`` as that many runs.

    Every property the comparison reads — the frame, the shape name, and each
    run's font, size and weight — is identical whatever the split, so two decks
    written by this function differ in the run count and in nothing else.
    """
    presentation = Presentation()
    presentation.slide_width = Inches(10)
    presentation.slide_height = Inches(5.625)
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])

    box = slide.shapes.add_textbox(Inches(1.00), Inches(1.00), Inches(6.0), Inches(0.5))
    box.name = "coalesced-body"
    paragraph = box.text_frame.paragraphs[0]
    for fragment in fragments:
        run = paragraph.add_run()
        run.text = fragment
        run.font.name = "Arial"
        run.font.size = Pt(18)
        run.font.bold = False

    presentation.save(str(path))


def main() -> int:
    FIXTURES.mkdir(parents=True, exist_ok=True)
    build(FIXTURES / "ref.pptx", title_size_pt=28, rect_x_in=RECT_X_IN, extra=False)
    build(FIXTURES / "gen.pptx", title_size_pt=24, rect_x_in=GEN_RECT_X_IN, extra=True)
    build_run_split(FIXTURES / "coalesced-ref.pptx", (COALESCED_TEXT,))
    build_run_split(FIXTURES / "coalesced-gen.pptx", SPLIT_TEXT)
    print(f"wrote ref, gen, coalesced-ref and coalesced-gen decks to {FIXTURES}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
