"""Tests for the pptx-compare structural comparison.

The fixtures under ``fixtures/`` are two pairs of one-slide decks regenerated
by ``gen_fixtures.py``. ``ref``/``gen`` differ by three planted discrepancies,
and each test asserts the parse finds exactly those and invents nothing else;
``coalesced-ref``/``coalesced-gen`` differ only by run coalescing, which is the
pair the ``--fold-engine-artefacts`` tests reduce to zero.

``compare_decks`` is imported as a plain module — putting its directory on
``sys.path`` is enough, because the scripts here are named with underscores.
The sibling pptx-arch-style tests need an ``importlib`` file loader instead,
because the scripts there are hyphenated and so cannot be imported by name.
That loader must not be copied into this tree.
"""

from __future__ import annotations

import sys
from pathlib import Path

import pytest

HERE = Path(__file__).resolve().parent
SCRIPTS = HERE.parent
FIXTURES = HERE / "fixtures"

if str(SCRIPTS) not in sys.path:
    sys.path.insert(0, str(SCRIPTS))

import compare_decks  # noqa: E402

REF = FIXTURES / "ref.pptx"
GEN = FIXTURES / "gen.pptx"
COALESCED_REF = FIXTURES / "coalesced-ref.pptx"
COALESCED_GEN = FIXTURES / "coalesced-gen.pptx"

# The 0.10in offset planted in gen.pptx, either side of which the coordinate
# discrepancy has to appear and disappear.
TOL_BELOW_OFFSET = compare_decks.DEFAULT_POS_TOL_EMU  # 0.04in
TOL_ABOVE_OFFSET = int(0.25 * compare_decks.EMU_PER_INCH)


def _lines(report) -> list[str]:
    return [line for slide in report.slides for line in slide.lines]


def test_fixtures_exist():
    missing = [
        deck.name
        for deck in (REF, GEN, COALESCED_REF, COALESCED_GEN)
        if not deck.exists()
    ]
    assert not missing, f"run gen_fixtures.py to build {missing}"


def test_deck_compared_with_itself_is_clean():
    report = compare_decks.compare_decks(REF, REF)
    assert report.ok
    assert report.diff_count == 0


def test_planted_discrepancies_are_all_found():
    report = compare_decks.compare_decks(REF, GEN, TOL_BELOW_OFFSET)
    assert not report.ok
    assert len(report.slides) == 1
    lines = _lines(report)
    assert len(lines) == 3, lines

    size, offset, unmatched = lines
    assert "size: ref=28.0 gen=24.0" in size
    assert "title" in size
    assert "x: ref=548640 gen=640080" in offset
    assert "body-box" in offset
    assert unmatched.startswith("only in gen:")
    assert "extra-note" in unmatched


def test_coordinate_delta_is_silenced_by_a_wider_tolerance():
    report = compare_decks.compare_decks(REF, GEN, TOL_ABOVE_OFFSET)
    lines = _lines(report)
    assert len(lines) == 2, lines
    assert not any("body-box" in line for line in lines)


def test_slide_count_mismatch_is_fatal(tmp_path):
    from pptx import Presentation

    short = tmp_path / "empty.pptx"
    Presentation().save(str(short))
    report = compare_decks.compare_decks(REF, short)
    assert report.fatal is not None
    assert "slide count differs" in report.fatal
    assert not report.ok


def test_matching_pairs_shapes_by_cost_not_by_index():
    """The extra gen shape must not shunt every later shape into a false diff."""
    ref_shapes = compare_decks.parse_deck(REF)[0]
    gen_shapes = compare_decks.parse_deck(GEN)[0]
    pairs = compare_decks.match_shapes(ref_shapes, gen_shapes)
    matched = {
        (a.name, b.name) for a, b in pairs if a is not None and b is not None
    }
    assert matched == {("title", "title"), ("body-box", "body-box")}
    assert [b.name for a, b in pairs if a is None] == ["extra-note"]


def test_cli_pos_tol_flag_is_exposed_in_inches():
    args = compare_decks.build_parser().parse_args([str(REF), str(GEN)])
    assert args.pos_tol == pytest.approx(
        compare_decks.DEFAULT_POS_TOL_EMU / compare_decks.EMU_PER_INCH
    )
    widened = compare_decks.build_parser().parse_args(
        [str(REF), str(GEN), "--pos-tol", "0.25"]
    )
    assert widened.pos_tol == pytest.approx(0.25)


def test_cli_outdir_has_no_default_beside_the_script():
    args = compare_decks.build_parser().parse_args([str(REF), str(GEN)])
    assert args.outdir is None


def test_exit_codes(capsys):
    assert compare_decks.main([str(REF), str(REF)]) == 0
    assert compare_decks.main([str(REF), str(GEN)]) == 1
    capsys.readouterr()


def test_render_without_soffice_fails_loudly(monkeypatch, tmp_path):
    """The missing-tool guard is all that stands between a user and a traceback."""
    monkeypatch.setattr(compare_decks.shutil, "which", lambda name: None)
    with pytest.raises(RuntimeError, match="soffice"):
        compare_decks.render_deck(REF, tmp_path / "out", 144)


def test_render_needs_pdftoppm_too(monkeypatch, tmp_path):
    monkeypatch.setattr(
        compare_decks.shutil, "which", lambda name: None if name == "pdftoppm" else "/x"
    )
    with pytest.raises(RuntimeError, match="pdftoppm"):
        compare_decks.render_deck(REF, tmp_path / "out", 144)


def test_render_failure_exits_2_and_writes_nothing(monkeypatch, tmp_path, capsys):
    monkeypatch.setattr(compare_decks.shutil, "which", lambda name: None)
    monkeypatch.chdir(tmp_path)
    assert compare_decks.main([str(REF), str(REF), "--render"]) == 2
    assert "render failed" in capsys.readouterr().err
    assert not (tmp_path / "_compare_out").exists()


def test_unreadable_deck_exits_2(tmp_path, capsys):
    junk = tmp_path / "junk.pptx"
    junk.write_text("not a zip archive at all", encoding="utf-8")
    assert compare_decks.main([str(REF), str(junk)]) == 2
    assert "not a readable .pptx" in capsys.readouterr().err


def test_missing_deck_exits_2(tmp_path, capsys):
    assert compare_decks.main([str(REF), str(tmp_path / "nope.pptx")]) == 2
    assert "no such deck" in capsys.readouterr().err


def test_format_report_marks_clean_slides_ok():
    text = compare_decks.format_report(compare_decks.compare_decks(REF, REF))
    assert "## Slide 1: OK" in text
    assert "Total: 0 discrepancies over 1 slide." in text


def test_format_report_lists_every_discrepancy():
    report = compare_decks.compare_decks(REF, GEN, TOL_BELOW_OFFSET)
    text = compare_decks.format_report(report)
    assert "## Slide 1: 3 discrepancies" in text
    assert text.count("\n- ") == 3 + 2  # three findings plus the ref/gen header pair
    assert "Total: 3 discrepancies over 1 slide." in text


def test_format_report_shows_a_fatal_instead_of_slides(tmp_path):
    from pptx import Presentation

    short = tmp_path / "empty.pptx"
    Presentation().save(str(short))
    text = compare_decks.format_report(compare_decks.compare_decks(REF, short))
    assert "FATAL: slide count differs" in text
    assert "## Slide" not in text


def test_report_flag_writes_the_same_text(tmp_path, capsys):
    destination = tmp_path / "nested" / "report.md"
    compare_decks.main([str(REF), str(GEN), "--report", str(destination)])
    written = destination.read_text(encoding="utf-8")
    assert written.strip() == capsys.readouterr().out.strip()


def test_groups_are_opaque_which_SKILL_md_documents(tmp_path):
    """Pin the limitation the skill warns about, so a future fix trips this test."""
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.util import Inches

    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    parts = [
        slide.shapes.add_shape(
            MSO_SHAPE.OVAL, Inches(1), Inches(1), Inches(1), Inches(1)
        ),
        slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(3), Inches(1), Inches(1), Inches(1)
        ),
    ]
    slide.shapes.add_group_shape(parts).name = "grp"
    deck = tmp_path / "grouped.pptx"
    presentation.save(str(deck))

    shapes = compare_decks.parse_deck(deck)[0]
    assert [s.name for s in shapes] == ["grp"], "children must not be walked"
    # The group must report no geometry of its own rather than borrowing a
    # child's preset, which would make the diff blame the wrong shape.
    assert shapes[0].geom is None


# --- --fold-engine-artefacts -------------------------------------------------
#
# The flag exists for the generator-convergence loop: fold what the engine pair
# explains, and "zero" becomes the single signal that the deck is aligned. Every
# test below is paired with its flag-off counterpart, because the whole point of
# the flag is that the default view stays exactly as honest as it was.


def test_cli_fold_flag_is_off_by_default():
    args = compare_decks.build_parser().parse_args([str(REF), str(GEN)])
    assert args.fold_engine_artefacts is False
    folded = compare_decks.build_parser().parse_args(
        [str(REF), str(GEN), "--fold-engine-artefacts"]
    )
    assert folded.fold_engine_artefacts is True


def test_cli_help_lists_the_fold_flag():
    help_text = compare_decks.build_parser().format_help()
    assert "--fold-engine-artefacts" in help_text


def test_run_coalescing_alone_is_reported_without_the_flag():
    """The default view must keep showing what the flag will later fold."""
    report = compare_decks.compare_decks(COALESCED_REF, COALESCED_GEN)
    assert not report.ok
    assert report.diff_count == 1
    (line,) = _lines(report)
    assert "run count: ref=1 gen=3" in line
    assert "engine artefact" in line


def test_run_coalescing_alone_folds_to_zero():
    report = compare_decks.compare_decks(
        COALESCED_REF, COALESCED_GEN, fold_engine_artefacts=True
    )
    assert report.ok
    assert report.diff_count == 0
    assert _lines(report) == []


def test_folded_report_totals_zero_and_exits_zero(capsys):
    """Zero is the convergence signal, in the total line and in the exit code."""
    assert compare_decks.main([str(COALESCED_REF), str(COALESCED_GEN)]) == 1
    capsys.readouterr()
    assert (
        compare_decks.main(
            [str(COALESCED_REF), str(COALESCED_GEN), "--fold-engine-artefacts"]
        )
        == 0
    )
    out = capsys.readouterr().out
    assert "Total: 0 discrepancies over 1 slide." in out
    assert "engine artefact if the text above matches" not in out


def test_folded_report_says_that_it_is_folded():
    """A filtered report has to admit it is filtered, or it misleads its reader."""
    folded = compare_decks.format_report(
        compare_decks.compare_decks(
            COALESCED_REF, COALESCED_GEN, fold_engine_artefacts=True
        )
    )
    assert "engine artefacts folded" in folded
    plain = compare_decks.format_report(
        compare_decks.compare_decks(COALESCED_REF, COALESCED_GEN)
    )
    assert "engine artefacts folded" not in plain


def test_flag_off_pins_the_prior_finding_count_on_the_prior_fixtures():
    """AC: the default output is unchanged — same three findings, no fold note."""
    report = compare_decks.compare_decks(REF, GEN, TOL_BELOW_OFFSET)
    assert report.diff_count == 3
    text = compare_decks.format_report(report)
    assert "Total: 3 discrepancies over 1 slide." in text
    assert "engine artefacts folded" not in text


def test_folding_leaves_real_drift_alone():
    """Folding must subtract artefacts only, never a genuine discrepancy."""
    report = compare_decks.compare_decks(
        REF, GEN, TOL_BELOW_OFFSET, fold_engine_artefacts=True
    )
    assert report.diff_count == 3
    assert _lines(report) == _lines(
        compare_decks.compare_decks(REF, GEN, TOL_BELOW_OFFSET)
    )


def _para(*fragments: str) -> compare_decks.Para:
    return compare_decks.Para(runs=[compare_decks.Run(text=f) for f in fragments])


def test_a_run_count_mismatch_survives_the_fold_when_the_text_differs():
    """The annotation hedges on "if the text above matches"; so does the fold.

    A differing run count next to differing text is not something the engine
    pair explains, so folding must leave it in place — otherwise the flag would
    hide drift, which is exactly what it is not for.
    """
    ref = _para("Alpha beta gamma")
    gen = _para("Alpha ", "beta ", "delta")
    kept = compare_decks.diff_runs(ref, gen, "para[0]", fold_engine_artefacts=True)
    assert any("text:" in line for line in kept)
    assert any("run count: ref=1 gen=3" in line for line in kept)


def test_folding_leaves_a_run_level_difference_alone():
    """Folding keys off the artefact, not off the paragraph having any runs.

    Equal run counts, one word bold on the gen side: nothing here is an
    artefact, so the fold must subtract nothing.
    """
    ref = compare_decks.Para(
        runs=[
            compare_decks.Run(text="Alpha ", bold=False),
            compare_decks.Run(text="beta", bold=False),
        ]
    )
    gen = compare_decks.Para(
        runs=[
            compare_decks.Run(text="Alpha ", bold=False),
            compare_decks.Run(text="beta", bold=True),
        ]
    )
    kept = compare_decks.diff_runs(ref, gen, "para[0]", fold_engine_artefacts=True)
    assert kept == ["para[0] run[1] bold: ref=False gen=True"]


def test_a_run_split_whose_formatting_differs_survives_the_fold():
    """Coalescing is lossless only for runs that carried the same formatting.

    The per-run loop zips to the shorter side, so with one run against three
    the formatting of gen's runs 1 and 2 is never inspected and the run-count
    line is the only trace they exist. Folding it on matching text alone would
    report a deck with a spuriously bold word as converged.
    """
    ref = compare_decks.Para(runs=[compare_decks.Run(text="Alpha beta gamma")])
    gen = compare_decks.Para(
        runs=[
            compare_decks.Run(text="Alpha "),
            compare_decks.Run(text="beta ", bold=True),
            compare_decks.Run(text="gamma"),
        ]
    )
    kept = compare_decks.diff_runs(ref, gen, "para[0]", fold_engine_artefacts=True)
    assert any("run count: ref=1 gen=3" in line for line in kept)


def test_a_mixed_format_paragraph_folds_when_the_merge_agrees():
    """The gate is coalescing-equivalence, not uniform formatting.

    A bold label followed by a plain value is a mixed paragraph, but splitting
    the value in two is still pure run splitting — so it still folds.
    """
    ref = compare_decks.Para(
        runs=[
            compare_decks.Run(text="Label: ", bold=True),
            compare_decks.Run(text="value"),
        ]
    )
    gen = compare_decks.Para(
        runs=[
            compare_decks.Run(text="Label: ", bold=True),
            compare_decks.Run(text="val"),
            compare_decks.Run(text="ue"),
        ]
    )
    assert (
        compare_decks.diff_runs(ref, gen, "para[0]", fold_engine_artefacts=True) == []
    )
