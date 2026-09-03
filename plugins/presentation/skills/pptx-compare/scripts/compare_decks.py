#!/usr/bin/env -S uv run --script
# /// script
# requires-python = ">=3.11"
# dependencies = [
#   "python-pptx>=1.0.2",
# ]
# ///
"""Structural comparison of two .pptx decks, shape by shape.

Parses both decks with python-pptx into a normalised shape model (text by run,
font face, size, weight/style, colour, fill, outline, alignment, indents, and
the x/y/w/h frame in EMU), matches the shapes of each slide pair by a pair
cost, then prints the discrepancies per slide.

The comparison is diagnostic and opinion-free: it measures the gap between two
decks and calls neither of them correct. Reading its output requires knowing
which discrepancies are engine artefacts rather than real drift — see
``references/engine-differences.md``.

Usage:
    uv run compare_decks.py REF.pptx GEN.pptx
    uv run compare_decks.py REF.pptx GEN.pptx --pos-tol 0.02
    uv run compare_decks.py REF.pptx GEN.pptx --render --dpi 144 --outdir ./_cmp
    uv run compare_decks.py REF.pptx GEN.pptx --report report.md

Exit codes:
    0 — the decks match within tolerance
    1 — at least one discrepancy
    2 — the decks could not be compared (slide count differs, render failed)
"""

from __future__ import annotations

import argparse
import difflib
import shutil
import subprocess
import sys
from dataclasses import dataclass, field
from pathlib import Path

from pptx import Presentation
from pptx.exc import PackageNotFoundError

EMU_PER_INCH = 914400
EMU_PER_POINT = 12700

# Default coordinate tolerance, ~0.04in. Exposed as --pos-tol because a deck
# built by a different generator needs a different slack.
DEFAULT_POS_TOL_EMU = 36576

# Weights of the pair-cost terms used to match a ref shape to a gen shape.
COST_TEXT_WEIGHT = 2.0
COST_GEOM_WEIGHT = 1.0
COST_KIND_PENALTY = 0.5
# Above this cost two shapes are considered unrelated and stay unmatched.
COST_MATCH_CEILING = 2.0


@dataclass
class Run:
    """One text run: the smallest span carrying uniform character formatting."""

    text: str
    font: str | None = None
    size_pt: float | None = None
    bold: bool | None = None
    italic: bool | None = None
    underline: bool | None = None
    color: str | None = None


@dataclass
class Para:
    """One paragraph: alignment, indents and its runs."""

    alignment: str | None = None
    level: int = 0
    space_before_pt: float | None = None
    space_after_pt: float | None = None
    indent_pt: float | None = None
    margin_left_pt: float | None = None
    runs: list[Run] = field(default_factory=list)

    @property
    def text(self) -> str:
        return "".join(r.text for r in self.runs)


@dataclass
class Shape:
    """A shape reduced to the properties worth diffing between two decks."""

    index: int
    name: str
    kind: str
    geom: str | None = None
    x: int | None = None
    y: int | None = None
    w: int | None = None
    h: int | None = None
    rotation: float = 0.0
    fill: str | None = None
    line: str | None = None
    line_width_pt: float | None = None
    word_wrap: bool | None = None
    paras: list[Para] = field(default_factory=list)

    @property
    def text(self) -> str:
        return "\n".join(p.text for p in self.paras)

    def label(self) -> str:
        head = self.text.strip().replace("\n", " / ")
        if len(head) > 48:
            head = head[:45] + "..."
        return f"[{self.index}] {self.kind} {self.name!r}" + (
            f" {head!r}" if head else ""
        )


@dataclass
class SlideDiff:
    """Discrepancies found on one slide pair."""

    number: int
    lines: list[str] = field(default_factory=list)

    @property
    def ok(self) -> bool:
        return not self.lines


@dataclass
class Report:
    """The result of comparing two decks."""

    ref: Path
    gen: Path
    slides: list[SlideDiff] = field(default_factory=list)
    fatal: str | None = None

    @property
    def ok(self) -> bool:
        return self.fatal is None and all(s.ok for s in self.slides)

    @property
    def diff_count(self) -> int:
        return sum(len(s.lines) for s in self.slides)


def _emu_to_pt(value: int | None) -> float | None:
    return None if value is None else round(value / EMU_PER_POINT, 2)


def _color_of(fmt) -> str | None:
    """Best-effort hex/scheme name of a colour, tolerant of unset formats."""
    try:
        color = fmt.color
    except (AttributeError, NotImplementedError):
        return None
    if color is None or color.type is None:
        return None
    try:
        return f"#{color.rgb}"
    except (AttributeError, TypeError, ValueError):
        pass
    try:
        return f"scheme:{color.theme_color}"
    except (AttributeError, TypeError, ValueError):
        return str(color.type)


def _fill_of(shape) -> str | None:
    try:
        fill = shape.fill
        kind = fill.type
    except (AttributeError, NotImplementedError, KeyError):
        return None
    if kind is None:
        return None
    name = str(kind).split(" ")[0].split(".")[-1]
    if "SOLID" in str(kind):
        return f"solid {_color_of(fill) or '?'}"
    return name.lower()


def _line_of(shape) -> tuple[str | None, float | None]:
    try:
        line = shape.line
    except (AttributeError, NotImplementedError, KeyError):
        return None, None
    width = _emu_to_pt(line.width) if line.width else None
    try:
        if line.fill.type is None:
            return None, width
    except (AttributeError, NotImplementedError, KeyError):
        return None, width
    return _color_of(line.fill) or "set", width


def _parse_runs(paragraph) -> list[Run]:
    out = []
    for run in paragraph.runs:
        font = run.font
        out.append(
            Run(
                text=run.text,
                font=font.name,
                size_pt=None if font.size is None else round(font.size.pt, 2),
                bold=font.bold,
                italic=font.italic,
                underline=font.underline,
                color=_color_of(font),
            )
        )
    return out


def _indents_of(paragraph) -> tuple[float | None, float | None]:
    """First-line indent and left margin in points, read off ``<a:pPr>``.

    python-pptx surfaces alignment and spacing on the paragraph but not
    ``marL``/``indent``, so those two come straight from the XML.
    """
    pPr = paragraph._p.find(
        "{http://schemas.openxmlformats.org/drawingml/2006/main}pPr"
    )
    if pPr is None:
        return None, None
    indent = pPr.get("indent")
    margin = pPr.get("marL")
    return (
        None if indent is None else round(int(indent) / EMU_PER_POINT, 2),
        None if margin is None else round(int(margin) / EMU_PER_POINT, 2),
    )


def _parse_paras(text_frame) -> list[Para]:
    out = []
    for paragraph in text_frame.paragraphs:
        indent_pt, margin_left_pt = _indents_of(paragraph)
        out.append(
            Para(
                alignment=(
                    None if paragraph.alignment is None else str(paragraph.alignment)
                ),
                level=paragraph.level,
                space_before_pt=(
                    None
                    if paragraph.space_before is None
                    else round(paragraph.space_before.pt, 2)
                ),
                space_after_pt=(
                    None
                    if paragraph.space_after is None
                    else round(paragraph.space_after.pt, 2)
                ),
                indent_pt=indent_pt,
                margin_left_pt=margin_left_pt,
                runs=_parse_runs(paragraph),
            )
        )
    return out


def _geom_of(shape) -> str | None:
    """Preset geometry name, or the connector kind for a ``p:cxnSp``.

    Read off the shape's own ``spPr`` rather than any descendant, so a group
    does not inherit the geometry of whichever child happens to come first.
    """
    try:
        element = shape._element
    except AttributeError:
        return None
    spPr = element.find(
        "{http://schemas.openxmlformats.org/presentationml/2006/main}spPr"
    )
    if spPr is None:
        spPr = element.find(
            "{http://schemas.openxmlformats.org/drawingml/2006/main}spPr"
        )
    if spPr is None:
        return None
    prst = spPr.find(
        "{http://schemas.openxmlformats.org/drawingml/2006/main}prstGeom"
    )
    return None if prst is None else prst.get("prst")


def parse_shape(shape, index: int) -> Shape:
    """Reduce a python-pptx shape to the diffable model."""
    line, line_width = _line_of(shape)
    parsed = Shape(
        index=index,
        name=shape.name or "",
        kind=str(shape.shape_type).split(" ")[0] if shape.shape_type else "UNKNOWN",
        geom=_geom_of(shape),
        x=shape.left,
        y=shape.top,
        w=shape.width,
        h=shape.height,
        rotation=getattr(shape, "rotation", 0.0) or 0.0,
        fill=_fill_of(shape),
        line=line,
        line_width_pt=line_width,
    )
    if shape.has_text_frame:
        parsed.word_wrap = shape.text_frame.word_wrap
        parsed.paras = _parse_paras(shape.text_frame)
    if getattr(shape, "has_table", False):
        for row_index, row in enumerate(shape.table.rows):
            for cell in row.cells:
                for para in _parse_paras(cell.text_frame):
                    para.level = row_index
                    parsed.paras.append(para)
    return parsed


def parse_deck(path: Path) -> list[list[Shape]]:
    """Parse a deck into a per-slide list of shapes."""
    presentation = Presentation(str(path))
    return [
        [parse_shape(shape, i) for i, shape in enumerate(slide.shapes)]
        for slide in presentation.slides
    ]


def _text_distance(a: str, b: str) -> float:
    if not a and not b:
        return 0.0
    return 1.0 - difflib.SequenceMatcher(None, a, b).ratio()


def _geom_distance(a: Shape, b: Shape) -> float:
    pairs = [(a.x, b.x), (a.y, b.y), (a.w, b.w), (a.h, b.h)]
    total = 0.0
    for left, right in pairs:
        if left is None or right is None:
            total += 1.0
        else:
            total += min(abs(left - right) / EMU_PER_INCH, 1.0)
    return total / len(pairs)


def pair_cost(a: Shape, b: Shape) -> float:
    """How unlike two shapes are; lower means a better match."""
    cost = COST_TEXT_WEIGHT * _text_distance(a.text, b.text)
    cost += COST_GEOM_WEIGHT * _geom_distance(a, b)
    if a.kind != b.kind:
        cost += COST_KIND_PENALTY
    return cost


def match_shapes(
    ref: list[Shape], gen: list[Shape]
) -> list[tuple[Shape | None, Shape | None]]:
    """Greedily pair ref shapes with gen shapes by ascending pair cost.

    Whatever stays above ``COST_MATCH_CEILING`` is reported as an unmatched
    shape on one side rather than force-fitted into a misleading diff.
    """
    candidates = sorted(
        (
            (pair_cost(a, b), i, j)
            for i, a in enumerate(ref)
            for j, b in enumerate(gen)
        ),
        key=lambda item: (item[0], item[1], item[2]),
    )
    used_ref: set[int] = set()
    used_gen: set[int] = set()
    pairs: list[tuple[int, int]] = []
    for cost, i, j in candidates:
        if cost > COST_MATCH_CEILING:
            break
        if i in used_ref or j in used_gen:
            continue
        used_ref.add(i)
        used_gen.add(j)
        pairs.append((i, j))

    matched: list[tuple[Shape | None, Shape | None]] = [
        (ref[i], gen[j]) for i, j in sorted(pairs)
    ]
    matched += [(shape, None) for i, shape in enumerate(ref) if i not in used_ref]
    matched += [(None, shape) for j, shape in enumerate(gen) if j not in used_gen]
    return matched


def _diff_field(out: list[str], what: str, ref, gen) -> None:
    if ref != gen:
        out.append(f"{what}: ref={ref!r} gen={gen!r}")


def diff_runs(ref: Para, gen: Para, prefix: str) -> list[str]:
    """Diff two paragraphs run by run.

    Run counts routinely differ between engines even when the rendered text is
    identical (PowerPoint coalesces adjacent runs on save), so a bare count
    mismatch is reported only when the concatenated text also differs.
    """
    out: list[str] = []
    if ref.text != gen.text:
        out.append(f"{prefix} text: ref={ref.text!r} gen={gen.text!r}")
    if len(ref.runs) != len(gen.runs):
        out.append(
            f"{prefix} run count: ref={len(ref.runs)} gen={len(gen.runs)} "
            "(engine artefact if the text above matches — see engine-differences.md)"
        )
    for i, (a, b) in enumerate(zip(ref.runs, gen.runs)):
        tag = f"{prefix} run[{i}]"
        _diff_field(out, f"{tag} font", a.font, b.font)
        _diff_field(out, f"{tag} size", a.size_pt, b.size_pt)
        _diff_field(out, f"{tag} bold", a.bold, b.bold)
        _diff_field(out, f"{tag} italic", a.italic, b.italic)
        _diff_field(out, f"{tag} underline", a.underline, b.underline)
        _diff_field(out, f"{tag} color", a.color, b.color)
    return out


def diff_shape(ref: Shape, gen: Shape, pos_tol_emu: int) -> list[str]:
    """All discrepancies between a matched pair of shapes."""
    out: list[str] = []
    for axis in ("x", "y", "w", "h"):
        a = getattr(ref, axis)
        b = getattr(gen, axis)
        if a is None or b is None:
            if a is not b:
                out.append(f"{axis}: ref={a} gen={b}")
            continue
        delta = abs(a - b)
        if delta > pos_tol_emu:
            out.append(
                f"{axis}: ref={a} gen={b} (delta {delta} EMU = "
                f"{delta / EMU_PER_INCH:.3f}in)"
            )
    _diff_field(out, "kind", ref.kind, gen.kind)
    _diff_field(out, "geom", ref.geom, gen.geom)
    _diff_field(out, "rotation", ref.rotation, gen.rotation)
    _diff_field(out, "fill", ref.fill, gen.fill)
    _diff_field(out, "line", ref.line, gen.line)
    _diff_field(out, "line width", ref.line_width_pt, gen.line_width_pt)
    _diff_field(out, "word wrap", ref.word_wrap, gen.word_wrap)

    if len(ref.paras) != len(gen.paras):
        out.append(f"paragraph count: ref={len(ref.paras)} gen={len(gen.paras)}")
    for i, (a, b) in enumerate(zip(ref.paras, gen.paras)):
        prefix = f"para[{i}]"
        _diff_field(out, f"{prefix} alignment", a.alignment, b.alignment)
        _diff_field(out, f"{prefix} level", a.level, b.level)
        _diff_field(out, f"{prefix} space before", a.space_before_pt, b.space_before_pt)
        _diff_field(out, f"{prefix} space after", a.space_after_pt, b.space_after_pt)
        _diff_field(out, f"{prefix} indent", a.indent_pt, b.indent_pt)
        _diff_field(out, f"{prefix} left margin", a.margin_left_pt, b.margin_left_pt)
        out += diff_runs(a, b, prefix)
    return out


def compare_decks(
    ref_path: Path, gen_path: Path, pos_tol_emu: int = DEFAULT_POS_TOL_EMU
) -> Report:
    """Compare two decks slide by slide and shape by shape."""
    report = Report(ref=ref_path, gen=gen_path)
    ref_slides = parse_deck(ref_path)
    gen_slides = parse_deck(gen_path)
    if len(ref_slides) != len(gen_slides):
        report.fatal = (
            f"slide count differs: ref={len(ref_slides)} gen={len(gen_slides)}"
        )
        return report

    for number, (ref_shapes, gen_shapes) in enumerate(
        zip(ref_slides, gen_slides), start=1
    ):
        slide = SlideDiff(number=number)
        for ref_shape, gen_shape in match_shapes(ref_shapes, gen_shapes):
            if ref_shape is None:
                slide.lines.append(f"only in gen: {gen_shape.label()}")
            elif gen_shape is None:
                slide.lines.append(f"only in ref: {ref_shape.label()}")
            else:
                for line in diff_shape(ref_shape, gen_shape, pos_tol_emu):
                    slide.lines.append(f"{ref_shape.label()} -> {line}")
        report.slides.append(slide)
    return report


def render_deck(deck: Path, outdir: Path, dpi: int) -> list[Path]:
    """Render a deck to one PNG per slide via ``soffice`` then ``pdftoppm``.

    Both decks must be rendered at the same dpi for ``pixel_diff.py`` to line
    the pages up. Raises ``RuntimeError`` when either tool is missing or fails.
    """
    for tool in ("soffice", "pdftoppm"):
        if shutil.which(tool) is None:
            raise RuntimeError(
                f"{tool} not found on PATH; --render needs soffice and pdftoppm"
            )
    outdir.mkdir(parents=True, exist_ok=True)
    subprocess.run(
        [
            "soffice",
            "--headless",
            "--convert-to",
            "pdf",
            "--outdir",
            str(outdir),
            str(deck),
        ],
        check=True,
        capture_output=True,
    )
    pdf = outdir / f"{deck.stem}.pdf"
    if not pdf.exists():
        raise RuntimeError(f"soffice produced no pdf for {deck}")
    subprocess.run(
        ["pdftoppm", "-png", "-r", str(dpi), str(pdf), str(outdir / deck.stem)],
        check=True,
        capture_output=True,
    )
    return sorted(outdir.glob(f"{deck.stem}-*.png"))


def format_report(report: Report) -> str:
    """Render the comparison as the text written to stdout and to --report."""
    lines = [
        "# Deck comparison",
        "",
        f"- ref: `{report.ref}`",
        f"- gen: `{report.gen}`",
        "",
    ]
    if report.fatal:
        lines += [f"FATAL: {report.fatal}", ""]
        return "\n".join(lines)

    for slide in report.slides:
        if slide.ok:
            lines.append(f"## Slide {slide.number}: OK")
            continue
        lines.append(f"## Slide {slide.number}: {len(slide.lines)} discrepancies")
        lines += [f"- {line}" for line in slide.lines]
        lines.append("")

    count = len(report.slides)
    lines += [
        "",
        f"Total: {report.diff_count} discrepancies over {count} "
        f"slide{'' if count == 1 else 's'}.",
    ]
    return "\n".join(lines)


def build_parser() -> argparse.ArgumentParser:
    parser = argparse.ArgumentParser(
        description="Compare two .pptx decks shape by shape.",
    )
    parser.add_argument("ref", type=Path, metavar="REF.pptx", help="reference deck")
    parser.add_argument("gen", type=Path, metavar="GEN.pptx", help="generated deck")
    parser.add_argument(
        "--pos-tol",
        type=float,
        default=DEFAULT_POS_TOL_EMU / EMU_PER_INCH,
        metavar="INCHES",
        help=(
            "coordinate tolerance in inches; x/y/w/h deltas at or below it are "
            "not reported (default: %(default).3f)"
        ),
    )
    parser.add_argument(
        "--render",
        action="store_true",
        help="also render both decks to PNG via soffice then pdftoppm",
    )
    parser.add_argument(
        "--dpi", type=int, default=144, help="render resolution (default: %(default)s)"
    )
    parser.add_argument(
        "--outdir",
        type=Path,
        default=None,
        metavar="DIR",
        help=(
            "where --render writes; defaults to ./_compare_out under the current "
            "working directory, never next to this script"
        ),
    )
    parser.add_argument(
        "--report", type=Path, default=None, metavar="FILE", help="also write here"
    )
    return parser


def main(argv: list[str] | None = None) -> int:
    args = build_parser().parse_args(argv)
    for deck in (args.ref, args.gen):
        if not deck.exists():
            print(f"no such deck: {deck}", file=sys.stderr)
            return 2

    pos_tol_emu = int(round(args.pos_tol * EMU_PER_INCH))
    try:
        report = compare_decks(args.ref, args.gen, pos_tol_emu)
    except PackageNotFoundError as exc:
        print(f"not a readable .pptx: {exc}", file=sys.stderr)
        return 2
    text = format_report(report)
    print(text)
    if args.report:
        args.report.parent.mkdir(parents=True, exist_ok=True)
        args.report.write_text(text + "\n", encoding="utf-8")

    if args.render:
        outdir = args.outdir or Path.cwd() / "_compare_out"
        try:
            for deck, sub in ((args.ref, "ref"), (args.gen, "gen")):
                pages = render_deck(deck, outdir / sub, args.dpi)
                print(f"rendered {len(pages)} page(s) of {deck.name} to {outdir / sub}")
        except (RuntimeError, subprocess.CalledProcessError) as exc:
            print(f"render failed: {exc}", file=sys.stderr)
            return 2

    if report.fatal:
        return 2
    return 0 if report.ok else 1


if __name__ == "__main__":
    raise SystemExit(main())
