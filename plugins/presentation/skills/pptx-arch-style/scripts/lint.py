#!/usr/bin/env python3
"""pptx-arch-style linter.

Reads a .pptx file with python-pptx, classifies every slide via a speaker-notes
tag (``<!--arch-style:content|title|section-->``), then evaluates rules from
``references/rules.yaml`` against the produced XML.

Usage:
    uv run plugins/presentation/skills/pptx-arch-style/scripts/lint.py deck.pptx
    uv run plugins/presentation/skills/pptx-arch-style/scripts/lint.py deck.pptx --json
    uv run plugins/presentation/skills/pptx-arch-style/scripts/lint.py deck.pptx --rules path/to/rules.yaml

Exit codes:
    0 — all checks passed (no errors, warnings allowed)
    1 — at least one ``severity: error`` violation OR untagged slide
    2 — only ``severity: warning`` violations (no errors)
"""

from __future__ import annotations

import argparse
import json
import re
import sys
from dataclasses import asdict, dataclass, field
from pathlib import Path

import yaml
from pptx import Presentation
from pptx.enum.text import PP_ALIGN

EMU_PER_INCH = 914400
EMU_PER_POINT = 12700
NS = {
    "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
    "p": "http://schemas.openxmlformats.org/presentationml/2006/main",
}
KIND_RE = re.compile(r"<!--\s*arch-style:(content|title|section)\s*-->")
DEFAULT_COORD_TOL = 0.005


@dataclass
class Violation:
    rule_id: str
    severity: str
    slide_num: int
    slide_kind: str
    message: str
    expected: str
    actual: str
    spec_ref: str


@dataclass
class Report:
    deck: str
    n_slides: int
    violations: list[Violation] = field(default_factory=list)
    untagged_slides: list[int] = field(default_factory=list)
    passed_checks: int = 0


def _emu_to_in(emu: int | None) -> float:
    return (emu or 0) / EMU_PER_INCH


def _emu_to_pt(emu: int | None) -> float:
    return (emu or 0) / EMU_PER_POINT


def slide_kind(slide) -> str | None:
    if not slide.has_notes_slide:
        return None
    text = slide.notes_slide.notes_text_frame.text or ""
    m = KIND_RE.search(text)
    return m.group(1) if m else None


def shape_fill_hex(shape) -> str | None:
    try:
        fill = shape.fill
    except (AttributeError, NotImplementedError):
        return None
    try:
        if fill.type != 1:
            return None
        rgb = fill.fore_color.rgb
        return str(rgb).upper() if rgb is not None else None
    except (AttributeError, ValueError, TypeError):
        return None


def shape_line_color_hex(shape) -> str | None:
    try:
        line = shape.line
        rgb = line.color.rgb
        return str(rgb).upper() if rgb is not None else None
    except (AttributeError, ValueError, TypeError, NotImplementedError, KeyError):
        return None


def shape_line_width_pt(shape) -> float | None:
    try:
        w = shape.line.width
        return None if w is None else _emu_to_pt(int(w))
    except (AttributeError, ValueError, TypeError):
        return None


def shape_max_font_size_pt(shape) -> float | None:
    if not shape.has_text_frame:
        return None
    sizes: list[float] = []
    for para in shape.text_frame.paragraphs:
        for run in para.runs:
            if run.font.size is not None:
                sizes.append(run.font.size.pt)
    return max(sizes) if sizes else None


def shape_text(shape) -> str:
    return shape.text_frame.text if shape.has_text_frame else ""


def shape_arrowheads(shape) -> set[str]:
    """Return the set of arrowhead ends present on a shape's line.

    Looks at <a:ln><a:headEnd type=".../> and <a:tailEnd type=".../> inside the
    shape's <p:spPr>. Returns a subset of {"begin", "end"}. An end is considered
    present iff its element exists AND its ``type`` attribute is anything other
    than ``none`` (default in OOXML is ``none``).
    """
    ends: set[str] = set()
    try:
        ln = shape._element.find("./p:spPr/a:ln", NS)
    except (AttributeError, ValueError, TypeError):
        return ends
    if ln is None:
        return ends
    for tag, key in (("a:headEnd", "begin"), ("a:tailEnd", "end")):
        el = ln.find(tag, NS)
        if el is not None and el.get("type", "none") != "none":
            ends.add(key)
    return ends


def shape_preset_geom(shape) -> str | None:
    """Return the prstGeom preset name for a shape, or None if absent.

    Handles auto-shapes (via ``auto_shape_type``), connector shapes
    (``shape_type``), and falls back to a direct XML lookup of
    ``a:prstGeom/@prst`` for shapes that python-pptx wraps as Shape/Picture
    without exposing a typed accessor.
    """
    try:
        ast = shape.auto_shape_type
        if ast is not None:
            return ast.name.upper()
    except (AttributeError, ValueError, TypeError):
        pass
    try:
        st = shape.shape_type
        if st is not None and st.name == "LINE":
            return "LINE"
    except (AttributeError, ValueError, TypeError):
        pass
    try:
        prst = shape._element.find(".//a:prstGeom", NS)
        if prst is not None:
            return prst.get("prst", "").upper() or None
    except (AttributeError, ValueError, TypeError):
        pass
    return None


def shape_matches(shape, match: dict | None) -> bool:
    if not match:
        return True

    fill = shape_fill_hex(shape)
    if "fill" in match:
        if fill is None or fill.upper() != match["fill"].upper():
            return False
    if "fill_in" in match:
        allowed = {c.upper() for c in match["fill_in"]}
        if fill is None or fill.upper() not in allowed:
            return False
    if "fill_not_in" in match:
        palette = {c.upper() for c in match["fill_not_in"]}
        if fill is None or fill.upper() in palette:
            return False

    x_in = _emu_to_in(shape.left)
    y_in = _emu_to_in(shape.top)
    w_in = _emu_to_in(shape.width)
    h_in = _emu_to_in(shape.height)
    bounds = {
        "x_min": (x_in, lambda a, b: a < b),
        "x_max": (x_in, lambda a, b: a > b),
        "y_min": (y_in, lambda a, b: a < b),
        "y_max": (y_in, lambda a, b: a > b),
        "w_min": (w_in, lambda a, b: a < b),
        "w_max": (w_in, lambda a, b: a > b),
        "h_min": (h_in, lambda a, b: a < b),
        "h_max": (h_in, lambda a, b: a > b),
    }
    for key, (actual, fail) in bounds.items():
        if key in match and fail(actual, match[key]):
            return False

    if "has_text" in match:
        has = shape.has_text_frame and bool(shape.text_frame.text.strip())
        if has != bool(match["has_text"]):
            return False

    if "font_size_min_pt" in match:
        mx = shape_max_font_size_pt(shape)
        if mx is None or mx < match["font_size_min_pt"]:
            return False

    if "line_color" in match:
        lc = shape_line_color_hex(shape)
        if lc is None or lc.upper() != match["line_color"].upper():
            return False

    if "shape_type" in match:
        expected = match["shape_type"].upper()
        actual = shape_preset_geom(shape)
        if actual != expected:
            return False

    if "arrowheads" in match:
        wanted = match["arrowheads"]
        ends = shape_arrowheads(shape)
        if wanted == "none" and ends:
            return False
        if wanted == "any" and not ends:
            return False
        if wanted == "begin" and "begin" not in ends:
            return False
        if wanted == "end" and "end" not in ends:
            return False

    return True


def bg_has_effectLst(slide) -> bool:
    elem = slide._element
    bg = elem.find(".//p:cSld/p:bg", NS)
    if bg is None:
        return False
    return bg.find(".//a:effectLst", NS) is not None


def _fmt_coord(d: dict) -> str:
    return ", ".join(f"{k}={v:.3f}" for k, v in d.items())


def _eval_mandatory(rule, slide_num, kind, slide, out):
    match = rule["expect"]["shape_match"]
    if any(shape_matches(s, match) for s in slide.shapes):
        return 1
    out.append(
        Violation(
            rule_id=rule["id"],
            severity=rule["severity"],
            slide_num=slide_num,
            slide_kind=kind,
            message=f"slide kind '{kind}' is missing a mandatory shape",
            expected=f"shape matching {match}",
            actual="not found",
            spec_ref=rule["spec_ref"],
        )
    )
    return 0


def _eval_forbidden(rule, slide_num, kind, slide, out):
    match = rule["forbid"]["shape_match"]
    found = [s for s in slide.shapes if shape_matches(s, match)]
    if not found:
        return 1
    out.append(
        Violation(
            rule_id=rule["id"],
            severity=rule["severity"],
            slide_num=slide_num,
            slide_kind=kind,
            message=f"slide kind '{kind}' contains forbidden shape",
            expected=f"no shape matching {match}",
            actual=f"{len(found)} match(es) found",
            spec_ref=rule["spec_ref"],
        )
    )
    return 0


def _eval_shape_coordinates(rule, slide_num, kind, slide, out):
    match = rule["applies_to"].get("shape_match")
    expect = rule["expect"]
    tol = rule.get("tolerance", {}).get("coord", DEFAULT_COORD_TOL)
    passed = 0
    for shape in slide.shapes:
        if not shape_matches(shape, match):
            continue
        actual = {
            "x": _emu_to_in(shape.left),
            "y": _emu_to_in(shape.top),
            "w": _emu_to_in(shape.width),
            "h": _emu_to_in(shape.height),
        }
        bad = {k: actual[k] for k in expect if abs(actual[k] - expect[k]) > tol}
        if bad:
            out.append(
                Violation(
                    rule_id=rule["id"],
                    severity=rule["severity"],
                    slide_num=slide_num,
                    slide_kind=kind,
                    message="shape coordinates outside tolerance",
                    expected=_fmt_coord({k: expect[k] for k in bad}),
                    actual=_fmt_coord(bad),
                    spec_ref=rule["spec_ref"],
                )
            )
        else:
            passed += 1
    return passed


def _eval_fill_color(rule, slide_num, kind, slide, out):
    match = rule["applies_to"].get("shape_match")
    msg = rule.get("message", "fill color violates palette")
    passed = 0
    for shape in slide.shapes:
        if shape_matches(shape, match):
            out.append(
                Violation(
                    rule_id=rule["id"],
                    severity=rule["severity"],
                    slide_num=slide_num,
                    slide_kind=kind,
                    message=msg,
                    expected="approved palette color",
                    actual=f"fill #{shape_fill_hex(shape)}",
                    spec_ref=rule["spec_ref"],
                )
            )
        else:
            passed += 1
    return passed


def _eval_border_spec(rule, slide_num, kind, slide, out):
    match = rule["applies_to"].get("shape_match")
    expect = rule["expect"]
    passed = 0
    for shape in slide.shapes:
        if not shape_matches(shape, match):
            continue
        color = shape_line_color_hex(shape)
        width_pt = shape_line_width_pt(shape)
        bad = []
        if "line_color" in expect and (color is None or color.upper() != expect["line_color"].upper()):
            bad.append(f"line_color={color} (expected {expect['line_color']})")
        if "line_width_pt_min" in expect and (width_pt is None or width_pt < expect["line_width_pt_min"]):
            bad.append(f"line_width_pt={width_pt} (min {expect['line_width_pt_min']})")
        if "line_width_pt_max" in expect and (width_pt is None or width_pt > expect["line_width_pt_max"]):
            bad.append(f"line_width_pt={width_pt} (max {expect['line_width_pt_max']})")
        if bad:
            out.append(
                Violation(
                    rule_id=rule["id"],
                    severity=rule["severity"],
                    slide_num=slide_num,
                    slide_kind=kind,
                    message="border spec mismatch",
                    expected=str(expect),
                    actual="; ".join(bad),
                    spec_ref=rule["spec_ref"],
                )
            )
        else:
            passed += 1
    return passed


def _eval_font_spec(rule, slide_num, kind, slide, out):
    faces = {f for f in rule["expect"]["faces"]}
    sizes = set(rule["expect"]["sizes_pt"])
    passed = 0
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                if not run.text or not run.text.strip():
                    continue
                face = run.font.name
                size = run.font.size.pt if run.font.size is not None else None
                bad = []
                if face and face not in faces:
                    bad.append(f"face='{face}'")
                if size is not None and size not in sizes:
                    bad.append(f"size={size}pt")
                if bad:
                    out.append(
                        Violation(
                            rule_id=rule["id"],
                            severity=rule["severity"],
                            slide_num=slide_num,
                            slide_kind=kind,
                            message=f"text run uses off-spec font: {run.text[:40]!r}",
                            expected=f"face in {sorted(faces)}; size in {sorted(sizes)}",
                            actual="; ".join(bad),
                            spec_ref=rule["spec_ref"],
                        )
                    )
                else:
                    passed += 1
    return passed


def _eval_text_alignment(rule, slide_num, kind, slide, out):
    match = rule["applies_to"].get("shape_match")
    align_map = {
        "left": PP_ALIGN.LEFT,
        "center": PP_ALIGN.CENTER,
        "right": PP_ALIGN.RIGHT,
    }
    expected_name = rule["expect"]["align"]
    target = align_map[expected_name]
    passed = 0
    for shape in slide.shapes:
        if not shape_matches(shape, match):
            continue
        if not shape.has_text_frame:
            continue
        for para in shape.text_frame.paragraphs:
            if para.alignment is None or para.alignment != target:
                out.append(
                    Violation(
                        rule_id=rule["id"],
                        severity=rule["severity"],
                        slide_num=slide_num,
                        slide_kind=kind,
                        message="text alignment mismatch",
                        expected=expected_name,
                        actual=str(para.alignment),
                        spec_ref=rule["spec_ref"],
                    )
                )
                break
        else:
            passed += 1
    return passed


def _eval_effect_override(rule, slide_num, kind, slide, out):
    if bg_has_effectLst(slide):
        return 1
    out.append(
        Violation(
            rule_id=rule["id"],
            severity=rule["severity"],
            slide_num=slide_num,
            slide_kind=kind,
            message="slide background lacks <a:effectLst/> shadow override",
            expected="<p:bgPr> contains <a:effectLst/>",
            actual="absent",
            spec_ref=rule["spec_ref"],
        )
    )
    return 0


EVALUATORS = {
    "mandatory_element": _eval_mandatory,
    "forbidden_element": _eval_forbidden,
    "shape_coordinates": _eval_shape_coordinates,
    "fill_color": _eval_fill_color,
    "border_spec": _eval_border_spec,
    "font_spec": _eval_font_spec,
    "text_alignment": _eval_text_alignment,
    "effect_override": _eval_effect_override,
}


def lint(deck_path: Path, rules: list[dict]) -> Report:
    deck = Presentation(str(deck_path))
    report = Report(deck=str(deck_path), n_slides=len(deck.slides))

    for idx, slide in enumerate(deck.slides, start=1):
        kind = slide_kind(slide)
        if kind is None:
            report.untagged_slides.append(idx)
            continue
        for rule in rules:
            applies = rule.get("applies_to", {})
            if "slide_kinds" in applies and kind not in applies["slide_kinds"]:
                continue
            passed = EVALUATORS[rule["type"]](rule, idx, kind, slide, report.violations)
            report.passed_checks += passed or 0

    return report


def format_text(report: Report) -> str:
    lines = [f"{report.deck} · {report.n_slides} slides", ""]
    by_slide: dict[int, list[Violation]] = {}
    for v in report.violations:
        by_slide.setdefault(v.slide_num, []).append(v)
    for idx in sorted({v.slide_num for v in report.violations}):
        vs = by_slide[idx]
        kind = vs[0].slide_kind.upper()
        lines.append(f"[Slide {idx} · {kind}]")
        for v in vs:
            marker = "x" if v.severity == "error" else "!"
            lines.append(f"  {marker} {v.rule_id} ({v.severity})")
            lines.append(f"      {v.message}")
            lines.append(f"      expected: {v.expected}")
            lines.append(f"      actual:   {v.actual}")
            lines.append(f"      spec:     {v.spec_ref}")
        lines.append("")
    for idx in report.untagged_slides:
        lines.append(f"[Slide {idx} · UNTAGGED]")
        lines.append("  x missing-classification-tag (error)")
        lines.append("      slide notes must contain <!--arch-style:content|title|section-->")
        lines.append("      spec:     SKILL.md → Validation gate (speaker-notes tagging)")
        lines.append("")
    errors = sum(1 for v in report.violations if v.severity == "error") + len(report.untagged_slides)
    warnings = sum(1 for v in report.violations if v.severity == "warning")
    lines.append(
        f"Summary: {report.passed_checks} passed · {errors} failed · {warnings} warnings"
    )
    return "\n".join(lines)


def exit_code(report: Report) -> int:
    errors = sum(1 for v in report.violations if v.severity == "error") + len(report.untagged_slides)
    if errors:
        return 1
    warnings = sum(1 for v in report.violations if v.severity == "warning")
    if warnings:
        return 2
    return 0


def load_rules(rules_path: Path) -> list[dict]:
    data = yaml.safe_load(rules_path.read_text(encoding="utf-8"))
    return data["rules"]


def main(argv: list[str] | None = None) -> int:
    ap = argparse.ArgumentParser(description="Lint a .pptx file against pptx-arch-style.")
    ap.add_argument("deck", type=Path, help="Path to .pptx file")
    ap.add_argument("--json", action="store_true", help="Emit machine-readable JSON")
    ap.add_argument(
        "--rules",
        type=Path,
        default=None,
        help="Override path to rules.yaml",
    )
    args = ap.parse_args(argv)

    rules_path = args.rules or (
        Path(__file__).resolve().parent.parent / "references" / "rules.yaml"
    )
    rules = load_rules(rules_path)
    report = lint(args.deck, rules)

    if args.json:
        payload = {
            "deck": report.deck,
            "n_slides": report.n_slides,
            "passed_checks": report.passed_checks,
            "untagged_slides": report.untagged_slides,
            "violations": [asdict(v) for v in report.violations],
        }
        print(json.dumps(payload, indent=2))
    else:
        print(format_text(report))

    code = exit_code(report)
    if args.json:
        sys.stderr.write(f"exit_code={code}\n")
    else:
        print(f"Exit code: {code}")
    return code


if __name__ == "__main__":
    sys.exit(main())
